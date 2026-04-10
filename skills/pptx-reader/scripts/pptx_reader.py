#!/usr/bin/env python3
"""
pptx-reader — extract the text content of a PowerPoint (.pptx) file.

One of the four readers in the `document-readers` plugin. Companion
to the `document-search` skill: use document-search to locate the
file, then this skill to read its content.

Usage:
    pptx_reader.py <file.pptx> [--slides 1-5,7]
                               [--max-bytes N]
                               [--format text|json]
                               [--metadata-only]

Run with --help for the full option list.
"""
from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Optional

sys.path.insert(0, str(Path(__file__).resolve().parent))
import _preflight  # noqa: E402


@dataclass
class ReadResult:
    path: str
    extension: str = ".pptx"
    kind: str = "pptx"
    size: int = 0
    content: Optional[str] = None
    truncated: bool = False
    metadata: dict = field(default_factory=dict)
    error: Optional[str] = None
    missing_imports: list[str] = field(default_factory=list)
    install_guide: Optional[str] = None


class MissingDependency(RuntimeError):
    def __init__(self, import_name: str, feature: str):
        self.import_name = import_name
        self.feature = feature
        super().__init__(f"missing optional package '{import_name}' for {feature}")


def _parse_slide_range(raw: str, total_slides: int) -> list[int]:
    """Parse '1-5', '3', '1,3-4,7' into a sorted list of 1-based slide numbers."""
    if not raw:
        return list(range(1, total_slides + 1))
    result: set[int] = set()
    for token in raw.split(","):
        token = token.strip()
        if not token:
            continue
        if "-" in token:
            a, _, b = token.partition("-")
            start, end = int(a), int(b)
            for n in range(max(1, start), min(total_slides, end) + 1):
                result.add(n)
        else:
            n = int(token)
            if 1 <= n <= total_slides:
                result.add(n)
    return sorted(result)


def read_pptx(path: Path, slides_spec: Optional[str]) -> tuple[str, dict]:
    try:
        from pptx import Presentation  # type: ignore  # python-pptx
    except ImportError as e:
        raise MissingDependency("pptx", ".pptx reading") from e
    prs = Presentation(str(path))
    total_slides = len(prs.slides)
    slides = _parse_slide_range(slides_spec or "", total_slides)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        if i not in slides:
            continue
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    metadata = {
        "total_slides": total_slides,
        "slides_returned": slides,
    }
    return "\n".join(parts), metadata


def read_file(
    path: Path,
    *,
    slides: Optional[str],
    max_bytes: int,
    metadata_only: bool,
) -> ReadResult:
    result = ReadResult(path=str(path), size=path.stat().st_size)
    try:
        content, meta = read_pptx(path, slides)
        result.metadata = meta
    except MissingDependency as e:
        result.error = str(e)
        result.missing_imports = [e.import_name]
        result.install_guide = _preflight.format_install_guide(
            [e.import_name],
            feature="reading .pptx files",
            skill_name="pptx-reader",
        )
        return result
    except Exception as e:  # noqa: BLE001
        result.error = f"read failed: {e}"
        return result

    if content is not None and len(content.encode("utf-8", errors="ignore")) > max_bytes:
        content = content.encode("utf-8", errors="ignore")[:max_bytes].decode("utf-8", errors="replace")
        result.truncated = True

    if not metadata_only:
        result.content = content
    return result


def render_text(r: ReadResult) -> str:
    out: list[str] = []
    out.append("# pptx-reader")
    out.append(f"path:      {r.path}")
    out.append(f"size:      {r.size} bytes")
    if r.truncated:
        out.append("truncated: yes (output capped by --max-bytes)")
    if r.metadata:
        out.append("metadata:")
        for k, v in r.metadata.items():
            out.append(f"  {k}: {v}")
    out.append("")

    if r.error:
        out.append(f"error: {r.error}")
        if r.install_guide:
            out.append("")
            out.append("## Install guide")
            out.append(r.install_guide)
        return "\n".join(out)

    if r.content is None:
        out.append("(metadata-only mode; no content returned)")
        return "\n".join(out)

    out.append("## Content")
    out.append(r.content)
    return "\n".join(out)


def render_json(r: ReadResult) -> str:
    return json.dumps(asdict(r), ensure_ascii=False, indent=2)


def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="pptx_reader.py",
        description="Extract the text content of a PowerPoint (.pptx) file, "
                    "optionally scoped to a slide range.",
    )
    p.add_argument("file", help="Path to the .pptx file to read.")
    p.add_argument("--slides", default=None,
                   help="Slide range (e.g. '1-5', '3', '1,3-4,7'). "
                        "Default: all slides.")
    p.add_argument("--max-bytes", type=int, default=200_000,
                   help="Maximum bytes of content returned (default: 200000).")
    p.add_argument("--format", choices=("text", "json"), default="text",
                   help="Output format (default: text).")
    p.add_argument("--metadata-only", action="store_true",
                   help="Return only metadata (slide count) without the content.")
    return p


def main(argv: Optional[list[str]] = None) -> int:
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")  # type: ignore[attr-defined]
        except (AttributeError, ValueError):
            pass

    args = build_parser().parse_args(argv)
    path = Path(args.file)
    if not path.exists():
        print(f"error: file does not exist: {path}", file=sys.stderr)
        return 2
    if not path.is_file():
        print(f"error: not a regular file: {path}", file=sys.stderr)
        return 2
    if path.suffix.lower() != ".pptx":
        print(f"warning: file does not have a .pptx extension: {path}", file=sys.stderr)

    result = read_file(
        path,
        slides=args.slides,
        max_bytes=args.max_bytes,
        metadata_only=args.metadata_only,
    )
    if args.format == "json":
        print(render_json(result))
    else:
        print(render_text(result))
    return 0 if result.error is None else 2


if __name__ == "__main__":
    sys.exit(main())
