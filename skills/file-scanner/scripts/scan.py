#!/usr/bin/env python3
"""
file-scanner — recursively scan a workspace and extract text from files.

Supports 60+ file types out of the box, including plain text / code files
and (optionally, if the relevant library is installed) Microsoft Office
documents and PDFs.

Designed to be invoked by AI coding agents (Claude Code, Roo Code, Aider,
Cursor, Cline, ...) to ground their context in actual workspace content.

Usage:
    scan.py <path> [--ext .py,.md] [--grep PATTERN] [--format text|json]
            [--max-bytes N] [--max-files N] [--ignore DIR,DIR]
            [--context N] [--list-only] [--include-hidden]

Run with --help for the full option list.
"""
from __future__ import annotations

import argparse
import fnmatch
import json
import os
import re
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Iterable, Iterator, Optional

# ---------------------------------------------------------------------------
# Supported extensions
# ---------------------------------------------------------------------------

# Plain-text / code files we can read directly with open().
TEXT_EXTENSIONS: set[str] = {
    # Programming languages
    ".py", ".pyi", ".pyx",
    ".js", ".jsx", ".ts", ".tsx", ".mjs", ".cjs",
    ".java", ".kt", ".kts", ".scala", ".groovy",
    ".c", ".h", ".cc", ".cpp", ".cxx", ".hpp", ".hh",
    ".cs", ".fs", ".vb",
    ".go", ".rs", ".swift", ".m", ".mm",
    ".rb", ".php", ".pl", ".pm", ".lua",
    ".r", ".jl", ".dart", ".ex", ".exs", ".erl", ".hrl",
    ".clj", ".cljs", ".cljc", ".edn",
    ".hs", ".lhs", ".ml", ".mli", ".nim", ".zig",
    # Shell / scripting
    ".sh", ".bash", ".zsh", ".fish", ".ps1", ".psm1", ".bat", ".cmd",
    # Markup / docs
    ".md", ".markdown", ".mdx", ".rst", ".txt", ".adoc", ".asciidoc",
    ".tex", ".org",
    # Data / config
    ".json", ".jsonc", ".json5", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".conf", ".properties", ".env",
    ".xml", ".html", ".htm", ".xhtml", ".svg",
    ".css", ".scss", ".sass", ".less",
    ".csv", ".tsv",
    # Build / infra
    ".dockerfile", ".tf", ".tfvars", ".hcl", ".nix",
    ".gradle", ".sbt", ".mk", ".cmake",
    # SQL
    ".sql",
    # Misc
    ".log", ".patch", ".diff",
}

# Binary / structured files we can read with an optional helper library.
BINARY_EXTENSIONS: set[str] = {".docx", ".pptx", ".xlsx", ".pdf"}

ALL_DEFAULT_EXTENSIONS: set[str] = TEXT_EXTENSIONS | BINARY_EXTENSIONS

# Directories we skip by default.
DEFAULT_IGNORES: set[str] = {
    ".git", ".hg", ".svn",
    "node_modules", "bower_components",
    "__pycache__", ".pytest_cache", ".mypy_cache", ".ruff_cache",
    "venv", ".venv", "env", ".env",
    "dist", "build", "target", "out",
    ".idea", ".vscode",
    ".next", ".nuxt", ".svelte-kit",
    ".gradle", ".tox",
}

# Filenames (no extension) that are still useful to scan as text.
EXTENSIONLESS_TEXT_NAMES: set[str] = {
    "Dockerfile", "Makefile", "Rakefile", "Gemfile", "Procfile",
    "Jenkinsfile", "Vagrantfile", "CMakeLists.txt",
    "README", "LICENSE", "CHANGELOG", "NOTICE", "AUTHORS",
}


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

@dataclass
class FileResult:
    path: str
    size: int
    extension: str
    kind: str                         # "text" | "docx" | "pptx" | "xlsx" | "pdf"
    truncated: bool = False
    content: Optional[str] = None     # None in --list-only mode
    matches: list[dict] = field(default_factory=list)  # for --grep
    error: Optional[str] = None


@dataclass
class ScanSummary:
    root: str
    files_scanned: int
    files_matched: int
    bytes_read: int
    skipped_missing_deps: dict[str, int]  # ext -> count
    results: list[FileResult]


# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------

def read_text_file(path: Path, max_bytes: int) -> tuple[str, bool]:
    """Read a text file up to max_bytes. Returns (content, truncated)."""
    with path.open("rb") as f:
        raw = f.read(max_bytes + 1)
    truncated = len(raw) > max_bytes
    if truncated:
        raw = raw[:max_bytes]
    # Decode leniently — the goal is to surface content to an LLM, not
    # to be byte-perfect.
    text = raw.decode("utf-8", errors="replace")
    return text, truncated


def read_docx(path: Path) -> str:
    try:
        import docx  # type: ignore  # python-docx
    except ImportError as e:
        raise RuntimeError("python-docx not installed") from e
    document = docx.Document(str(path))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore  # python-pptx
    except ImportError as e:
        raise RuntimeError("python-pptx not installed") from e
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def read_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError as e:
        raise RuntimeError("openpyxl not installed") from e
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def read_pdf(path: Path) -> str:
    try:
        import pypdf  # type: ignore
    except ImportError as e:
        raise RuntimeError("pypdf not installed") from e
    reader = pypdf.PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            text = f"<error extracting page {i}: {exc}>"
        parts.append(f"--- Page {i} ---")
        parts.append(text)
    return "\n".join(parts)


BINARY_READERS = {
    ".docx": ("docx", read_docx),
    ".pptx": ("pptx", read_pptx),
    ".xlsx": ("xlsx", read_xlsx),
    ".pdf":  ("pdf",  read_pdf),
}


# ---------------------------------------------------------------------------
# Walking
# ---------------------------------------------------------------------------

def iter_files(
    root: Path,
    extensions: set[str],
    ignore_dirs: set[str],
    include_hidden: bool,
) -> Iterator[Path]:
    """Yield files under `root` that match the requested extensions."""
    for dirpath, dirnames, filenames in os.walk(root):
        # Prune ignored / hidden directories in-place so we don't descend.
        pruned = []
        for d in list(dirnames):
            if d in ignore_dirs:
                continue
            if not include_hidden and d.startswith("."):
                continue
            pruned.append(d)
        dirnames[:] = pruned

        # Did the caller ask for any "text" extension? If so we also surface
        # extension-less files like Makefile / Dockerfile.
        text_requested = bool(extensions & TEXT_EXTENSIONS)

        for name in filenames:
            if not include_hidden and name.startswith("."):
                continue
            full = Path(dirpath) / name
            ext = full.suffix.lower()
            if ext in extensions:
                yield full
            elif not ext and text_requested and name in EXTENSIONLESS_TEXT_NAMES:
                yield full


# ---------------------------------------------------------------------------
# Grep
# ---------------------------------------------------------------------------

def find_matches(text: str, pattern: re.Pattern, context: int) -> list[dict]:
    """Return a list of {line, text, before, after} match dicts."""
    if context < 0:
        context = 0
    lines = text.splitlines()
    matches: list[dict] = []
    for idx, line in enumerate(lines):
        if pattern.search(line):
            start = max(0, idx - context)
            end = min(len(lines), idx + context + 1)
            matches.append({
                "line": idx + 1,
                "text": line,
                "before": lines[start:idx],
                "after":  lines[idx + 1:end],
            })
    return matches


# ---------------------------------------------------------------------------
# Core scan
# ---------------------------------------------------------------------------

def scan(
    root: Path,
    *,
    extensions: set[str],
    ignore_dirs: set[str],
    grep: Optional[re.Pattern],
    max_bytes: int,
    max_files: int,
    include_hidden: bool,
    list_only: bool,
    context: int,
) -> ScanSummary:
    summary = ScanSummary(
        root=str(root.resolve()),
        files_scanned=0,
        files_matched=0,
        bytes_read=0,
        skipped_missing_deps={},
        results=[],
    )

    count = 0
    for path in iter_files(root, extensions, ignore_dirs, include_hidden):
        if max_files and count >= max_files:
            break
        count += 1

        try:
            size = path.stat().st_size
        except OSError as e:
            summary.results.append(FileResult(
                path=str(path), size=0, extension=path.suffix.lower(),
                kind="text", error=f"stat failed: {e}",
            ))
            continue

        ext = path.suffix.lower()
        result = FileResult(
            path=str(path),
            size=size,
            extension=ext,
            kind="text",
        )

        if list_only and not grep:
            summary.files_scanned += 1
            summary.results.append(result)
            continue

        # Extract content
        content: Optional[str] = None
        try:
            if ext in BINARY_READERS:
                kind, reader = BINARY_READERS[ext]
                result.kind = kind
                content = reader(path)
                if len(content.encode("utf-8", errors="ignore")) > max_bytes:
                    content = content.encode("utf-8", errors="ignore")[:max_bytes].decode("utf-8", errors="replace")
                    result.truncated = True
            else:
                content, result.truncated = read_text_file(path, max_bytes)
        except RuntimeError as e:
            # Missing optional dep
            result.error = str(e)
            summary.skipped_missing_deps[ext] = summary.skipped_missing_deps.get(ext, 0) + 1
            summary.results.append(result)
            summary.files_scanned += 1
            continue
        except Exception as e:  # noqa: BLE001
            result.error = f"read failed: {e}"
            summary.results.append(result)
            summary.files_scanned += 1
            continue

        summary.bytes_read += len(content.encode("utf-8", errors="ignore"))

        if grep is not None:
            matches = find_matches(content, grep, context)
            if not matches:
                summary.files_scanned += 1
                continue
            result.matches = matches
            summary.files_matched += 1

        if not list_only:
            result.content = content

        summary.results.append(result)
        summary.files_scanned += 1

    return summary


# ---------------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------------

def render_text(summary: ScanSummary, show_content: bool) -> str:
    out: list[str] = []
    out.append(f"# file-scanner results")
    out.append(f"root: {summary.root}")
    out.append(f"files scanned: {summary.files_scanned}")
    if summary.files_matched:
        out.append(f"files matched: {summary.files_matched}")
    out.append(f"bytes read:    {summary.bytes_read}")
    if summary.skipped_missing_deps:
        deps = ", ".join(f"{k} ({v})" for k, v in summary.skipped_missing_deps.items())
        out.append(f"skipped (missing optional deps): {deps}")
    out.append("")

    for r in summary.results:
        header = f"=== {r.path} [{r.kind}, {r.size} bytes"
        if r.truncated:
            header += ", truncated"
        header += "] ==="
        out.append(header)
        if r.error:
            out.append(f"  error: {r.error}")
            out.append("")
            continue
        if r.matches:
            for m in r.matches:
                for b in m["before"]:
                    out.append(f"  {m['line'] - len(m['before'])}- {b}")
                out.append(f"  {m['line']}: {m['text']}")
                for i, a in enumerate(m["after"], start=1):
                    out.append(f"  {m['line'] + i}- {a}")
                out.append("")
        elif show_content and r.content is not None:
            out.append(r.content)
            out.append("")
        else:
            out.append("")
    return "\n".join(out)


def render_json(summary: ScanSummary) -> str:
    payload = {
        "root": summary.root,
        "files_scanned": summary.files_scanned,
        "files_matched": summary.files_matched,
        "bytes_read": summary.bytes_read,
        "skipped_missing_deps": summary.skipped_missing_deps,
        "results": [asdict(r) for r in summary.results],
    }
    return json.dumps(payload, ensure_ascii=False, indent=2)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def parse_extensions(raw: Optional[str]) -> set[str]:
    if not raw:
        return set(ALL_DEFAULT_EXTENSIONS)
    result: set[str] = set()
    for token in raw.split(","):
        token = token.strip().lower()
        if not token:
            continue
        if not token.startswith("."):
            token = "." + token
        result.add(token)
    return result


def parse_ignores(raw: Optional[str]) -> set[str]:
    result = set(DEFAULT_IGNORES)
    if raw:
        for token in raw.split(","):
            token = token.strip()
            if token:
                result.add(token)
    return result


def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="scan.py",
        description="Recursively scan a workspace and extract text from files. "
                    "Supports 60+ file types including .py, .md, .docx, .pdf, .pptx, .xlsx.",
    )
    p.add_argument("path", nargs="?", default=".",
                   help="Root directory to scan (default: current directory).")
    p.add_argument("--ext", default=None,
                   help="Comma-separated list of extensions to include "
                        "(e.g. '.py,.md'). Defaults to 60+ built-in types.")
    p.add_argument("--grep", default=None,
                   help="Regex pattern to search for. Only files with matches "
                        "will be returned.")
    p.add_argument("--ignore-case", action="store_true",
                   help="Case-insensitive grep.")
    p.add_argument("--context", type=int, default=0,
                   help="Lines of context around each grep match (default: 0).")
    p.add_argument("--max-bytes", type=int, default=200_000,
                   help="Maximum bytes to read per file (default: 200000).")
    p.add_argument("--max-files", type=int, default=0,
                   help="Stop after scanning this many files (0 = no limit).")
    p.add_argument("--ignore", default=None,
                   help="Additional comma-separated directory names to skip.")
    p.add_argument("--include-hidden", action="store_true",
                   help="Include dotfiles and dot-directories.")
    p.add_argument("--list-only", action="store_true",
                   help="List matching files without emitting their content.")
    p.add_argument("--format", choices=("text", "json"), default="text",
                   help="Output format (default: text).")
    return p


def main(argv: Optional[list[str]] = None) -> int:
    args = build_parser().parse_args(argv)

    root = Path(args.path)
    if not root.exists():
        print(f"error: path does not exist: {root}", file=sys.stderr)
        return 2
    if not root.is_dir():
        print(f"error: path is not a directory: {root}", file=sys.stderr)
        return 2

    extensions = parse_extensions(args.ext)
    ignore_dirs = parse_ignores(args.ignore)

    grep_pat: Optional[re.Pattern] = None
    if args.grep:
        flags = re.IGNORECASE if args.ignore_case else 0
        try:
            grep_pat = re.compile(args.grep, flags)
        except re.error as e:
            print(f"error: invalid regex: {e}", file=sys.stderr)
            return 2

    summary = scan(
        root,
        extensions=extensions,
        ignore_dirs=ignore_dirs,
        grep=grep_pat,
        max_bytes=args.max_bytes,
        max_files=args.max_files,
        include_hidden=args.include_hidden,
        list_only=args.list_only,
        context=args.context,
    )

    if args.format == "json":
        print(render_json(summary))
    else:
        show_content = not args.list_only and grep_pat is None
        print(render_text(summary, show_content=show_content))
    return 0


if __name__ == "__main__":
    sys.exit(main())
