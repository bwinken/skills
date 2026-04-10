#!/usr/bin/env python3
"""
document-search — grep for a term across a workspace, even inside Office and PDF files.

Unlike ordinary `grep` / `rg`, this skill can look inside binary-wrapped
document formats (.docx, .pptx, .xlsx, .pdf) as well as 60+ text and code
file types. The default output is a compact list of files that contain the
search term, sorted by match count — exactly what an AI agent needs when a
user asks "which files mention X?".

Usage:
    document_search.py <pattern> [path] [--ext .py,.md,.docx]
                 [--fixed-string] [--ignore-case]
                 [--show-matches] [--context N]
                 [--max-bytes N] [--max-files N]
                 [--ignore DIR,DIR] [--include-hidden]
                 [--format text|json]

Run with --help for the full option list.
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Iterator, Optional

# Local, self-contained helper for install-guide generation. Kept as a
# sibling so the skill folder stays copy-and-go.
sys.path.insert(0, str(Path(__file__).resolve().parent))
import _preflight  # noqa: E402

# ---------------------------------------------------------------------------
# Supported extensions
# ---------------------------------------------------------------------------

# Plain-text / code files we can read directly with open().
TEXT_EXTENSIONS: set[str] = {
    # Programming languages
    ".py", ".pyi", ".pyx",
    ".js", ".jsx", ".ts", ".tsx", ".mjs", ".cjs",
    ".java", ".kt", ".kts", ".scala", ".groovy",
    ".c", ".h", ".cc", ".cpp", ".cxx", ".hpp", ".hh",
    ".cs", ".fs", ".vb",
    ".go", ".rs", ".swift", ".m", ".mm",
    ".rb", ".php", ".pl", ".pm", ".lua",
    ".r", ".jl", ".dart", ".ex", ".exs", ".erl", ".hrl",
    ".clj", ".cljs", ".cljc", ".edn",
    ".hs", ".lhs", ".ml", ".mli", ".nim", ".zig",
    # Shell / scripting
    ".sh", ".bash", ".zsh", ".fish", ".ps1", ".psm1", ".bat", ".cmd",
    # Markup / docs
    ".md", ".markdown", ".mdx", ".rst", ".txt", ".adoc", ".asciidoc",
    ".tex", ".org",
    # Data / config
    ".json", ".jsonc", ".json5", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".conf", ".properties", ".env",
    ".xml", ".html", ".htm", ".xhtml", ".svg",
    ".css", ".scss", ".sass", ".less",
    ".csv", ".tsv",
    # Build / infra
    ".dockerfile", ".tf", ".tfvars", ".hcl", ".nix",
    ".gradle", ".sbt", ".mk", ".cmake",
    # SQL
    ".sql",
    # Misc
    ".log", ".patch", ".diff",
}

# Binary / structured files we can read with an optional helper library.
BINARY_EXTENSIONS: set[str] = {".docx", ".pptx", ".xlsx", ".pdf"}

ALL_DEFAULT_EXTENSIONS: set[str] = TEXT_EXTENSIONS | BINARY_EXTENSIONS

# Directories we skip by default.
DEFAULT_IGNORES: set[str] = {
    ".git", ".hg", ".svn",
    "node_modules", "bower_components",
    "__pycache__", ".pytest_cache", ".mypy_cache", ".ruff_cache",
    "venv", ".venv", "env", ".env",
    "dist", "build", "target", "out",
    ".idea", ".vscode",
    ".next", ".nuxt", ".svelte-kit",
    ".gradle", ".tox",
}

# Filenames (no extension) that are still useful to scan as text.
EXTENSIONLESS_TEXT_NAMES: set[str] = {
    "Dockerfile", "Makefile", "Rakefile", "Gemfile", "Procfile",
    "Jenkinsfile", "Vagrantfile", "CMakeLists.txt",
    "README", "LICENSE", "CHANGELOG", "NOTICE", "AUTHORS",
}


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

@dataclass
class Match:
    line: int
    text: str
    before: list[str] = field(default_factory=list)
    after: list[str] = field(default_factory=list)


@dataclass
class FileHit:
    path: str
    extension: str
    kind: str                         # "text" | "docx" | "pptx" | "xlsx" | "pdf"
    size: int
    match_count: int
    truncated: bool = False
    matches: list[Match] = field(default_factory=list)  # empty unless --show-matches
    content: Optional[str] = None     # populated only with --content
    error: Optional[str] = None


@dataclass
class ScanSummary:
    pattern: str
    root: str
    files_scanned: int
    files_matched: int
    total_matches: int
    skipped_missing_deps: dict[str, int]         # ext -> count
    missing_imports: list[str] = field(default_factory=list)  # e.g. ["docx", "pypdf"]
    install_guide: Optional[str] = None          # populated only if missing_imports
    results: list[FileHit] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------

def read_text_file(path: Path, max_bytes: int) -> tuple[str, bool]:
    """Read a text file up to max_bytes. Returns (content, truncated)."""
    with path.open("rb") as f:
        raw = f.read(max_bytes + 1)
    truncated = len(raw) > max_bytes
    if truncated:
        raw = raw[:max_bytes]
    return raw.decode("utf-8", errors="replace"), truncated


class MissingDependency(RuntimeError):
    """Raised when an optional reader library is not installed."""
    def __init__(self, import_name: str, feature: str):
        self.import_name = import_name
        self.feature = feature
        super().__init__(f"missing optional package '{import_name}' for {feature}")


def read_docx(path: Path) -> str:
    try:
        import docx  # type: ignore  # python-docx
    except ImportError as e:
        raise MissingDependency("docx", ".docx reading") from e
    document = docx.Document(str(path))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore  # python-pptx
    except ImportError as e:
        raise MissingDependency("pptx", ".pptx reading") from e
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def read_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError as e:
        raise MissingDependency("openpyxl", ".xlsx reading") from e
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def read_pdf(path: Path) -> str:
    try:
        import pypdf  # type: ignore
    except ImportError as e:
        raise MissingDependency("pypdf", ".pdf reading") from e
    reader = pypdf.PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            text = f"<error extracting page {i}: {exc}>"
        parts.append(f"--- Page {i} ---")
        parts.append(text)
    return "\n".join(parts)


BINARY_READERS = {
    ".docx": ("docx", read_docx),
    ".pptx": ("pptx", read_pptx),
    ".xlsx": ("xlsx", read_xlsx),
    ".pdf":  ("pdf",  read_pdf),
}


# ---------------------------------------------------------------------------
# Walking
# ---------------------------------------------------------------------------

def iter_files(
    root: Path,
    extensions: set[str],
    ignore_dirs: set[str],
    include_hidden: bool,
) -> Iterator[Path]:
    """Yield files under `root` that match the requested extensions."""
    for dirpath, dirnames, filenames in os.walk(root):
        # Prune ignored / hidden directories in-place so we don't descend.
        pruned = []
        for d in list(dirnames):
            if d in ignore_dirs:
                continue
            if not include_hidden and d.startswith("."):
                continue
            pruned.append(d)
        dirnames[:] = pruned

        # Did the caller ask for any "text" extension? If so we also surface
        # extension-less files like Makefile / Dockerfile.
        text_requested = bool(extensions & TEXT_EXTENSIONS)

        for name in filenames:
            if not include_hidden and name.startswith("."):
                continue
            full = Path(dirpath) / name
            ext = full.suffix.lower()
            if ext in extensions:
                yield full
            elif not ext and text_requested and name in EXTENSIONLESS_TEXT_NAMES:
                yield full


# ---------------------------------------------------------------------------
# Match logic
# ---------------------------------------------------------------------------

def count_and_collect(
    text: str,
    pattern: re.Pattern,
    *,
    collect: bool,
    context: int,
) -> tuple[int, list[Match]]:
    """
    Return (match_count, matches). If `collect` is False, `matches` is empty
    but match_count is still accurate.
    """
    if not collect:
        return len(pattern.findall(text)), []

    lines = text.splitlines()
    matches: list[Match] = []
    total = 0
    for idx, line in enumerate(lines):
        hits = len(pattern.findall(line))
        if hits:
            total += hits
            start = max(0, idx - context)
            end = min(len(lines), idx + context + 1)
            matches.append(Match(
                line=idx + 1,
                text=line,
                before=lines[start:idx],
                after=lines[idx + 1:end],
            ))
    return total, matches


# ---------------------------------------------------------------------------
# Core scan
# ---------------------------------------------------------------------------

def scan(
    root: Path,
    pattern: re.Pattern,
    *,
    extensions: set[str],
    ignore_dirs: set[str],
    max_bytes: int,
    max_files: int,
    include_hidden: bool,
    show_matches: bool,
    include_content: bool,
    context: int,
) -> ScanSummary:
    summary = ScanSummary(
        pattern=pattern.pattern,
        root=str(root.resolve()),
        files_scanned=0,
        files_matched=0,
        total_matches=0,
        skipped_missing_deps={},
        results=[],
    )
    missing_imports_set: set[str] = set()

    count = 0
    for path in iter_files(root, extensions, ignore_dirs, include_hidden):
        if max_files and count >= max_files:
            break
        count += 1
        summary.files_scanned += 1

        try:
            size = path.stat().st_size
        except OSError as e:
            summary.results.append(FileHit(
                path=str(path), extension=path.suffix.lower(), kind="text",
                size=0, match_count=0, error=f"stat failed: {e}",
            ))
            continue

        ext = path.suffix.lower()
        kind = "text"
        content: Optional[str] = None
        truncated = False

        try:
            if ext in BINARY_READERS:
                kind, reader = BINARY_READERS[ext]
                content = reader(path)
                encoded = content.encode("utf-8", errors="ignore")
                if len(encoded) > max_bytes:
                    content = encoded[:max_bytes].decode("utf-8", errors="replace")
                    truncated = True
            else:
                content, truncated = read_text_file(path, max_bytes)
        except MissingDependency as e:
            # Missing optional dep — report but keep going.
            summary.skipped_missing_deps[ext] = summary.skipped_missing_deps.get(ext, 0) + 1
            missing_imports_set.add(e.import_name)
            summary.results.append(FileHit(
                path=str(path), extension=ext, kind=kind, size=size,
                match_count=0, error=str(e),
            ))
            continue
        except Exception as e:  # noqa: BLE001
            summary.results.append(FileHit(
                path=str(path), extension=ext, kind=kind, size=size,
                match_count=0, error=f"read failed: {e}",
            ))
            continue

        match_count, matches = count_and_collect(
            content, pattern, collect=show_matches, context=context,
        )
        if match_count == 0:
            continue

        summary.files_matched += 1
        summary.total_matches += match_count
        summary.results.append(FileHit(
            path=str(path),
            extension=ext,
            kind=kind,
            size=size,
            match_count=match_count,
            truncated=truncated,
            matches=matches,
            content=content if include_content else None,
        ))

    # Sort: successful hits first (by match count desc, then path), errored
    # entries at the end. The renderer splits them into separate sections.
    summary.results.sort(key=lambda r: (r.error is not None, -r.match_count, r.path))

    if missing_imports_set:
        summary.missing_imports = sorted(missing_imports_set)
        summary.install_guide = _preflight.format_install_guide(
            summary.missing_imports,
            feature="reading .docx / .pptx / .xlsx / .pdf files",
            skill_name="document-search",
        )
    return summary


# ---------------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------------

def render_text(summary: ScanSummary, show_matches: bool) -> str:
    out: list[str] = []
    out.append("# document-search results")
    out.append(f"pattern: {summary.pattern!r}")
    out.append(f"root:    {summary.root}")
    out.append(f"files matched: {summary.files_matched} / {summary.files_scanned} scanned "
               f"({summary.total_matches} total matches)")
    if summary.skipped_missing_deps:
        deps = ", ".join(f"{k} ({v})" for k, v in summary.skipped_missing_deps.items())
        out.append(f"skipped (missing optional deps): {deps}")
    out.append("")

    hits = [r for r in summary.results if not r.error]
    errored = [r for r in summary.results if r.error]

    if not hits and not errored:
        out.append("(no matches)")
        return "\n".join(out)

    if hits:
        out.append("## Matched files")
        path_width = max(len(r.path) for r in hits)
        for r in hits:
            tag = f" [{r.kind}]" if r.kind != "text" else ""
            noun = "match" if r.match_count == 1 else "matches"
            trunc = " (truncated)" if r.truncated else ""
            out.append(f"  {r.path.ljust(path_width)}  {r.match_count:>4} {noun}{tag}{trunc}")

    if errored:
        out.append("")
        out.append("## Skipped / errored")
        for r in errored:
            out.append(f"  {r.path}  — {r.error}")

    if summary.install_guide:
        out.append("")
        out.append("## Install guide for skipped files")
        out.append(summary.install_guide)

    if show_matches and hits:
        out.append("")
        out.append("## Match details")
        for r in hits:
            out.append(f"=== {r.path} ({r.match_count}) ===")
            for m in r.matches:
                for i, b in enumerate(m.before):
                    ln = m.line - len(m.before) + i
                    out.append(f"  {ln:>5}- {b}")
                out.append(f"  {m.line:>5}: {m.text}")
                for i, a in enumerate(m.after, start=1):
                    out.append(f"  {m.line + i:>5}- {a}")
                out.append("")
    return "\n".join(out)


def render_json(summary: ScanSummary) -> str:
    payload = {
        "pattern": summary.pattern,
        "root": summary.root,
        "files_scanned": summary.files_scanned,
        "files_matched": summary.files_matched,
        "total_matches": summary.total_matches,
        "skipped_missing_deps": summary.skipped_missing_deps,
        "missing_imports": summary.missing_imports,
        "install_guide": summary.install_guide,
        "results": [asdict(r) for r in summary.results],
    }
    return json.dumps(payload, ensure_ascii=False, indent=2)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def parse_extensions(raw: Optional[str]) -> set[str]:
    if not raw:
        return set(ALL_DEFAULT_EXTENSIONS)
    result: set[str] = set()
    for token in raw.split(","):
        token = token.strip().lower()
        if not token:
            continue
        if not token.startswith("."):
            token = "." + token
        result.add(token)
    return result


def parse_ignores(raw: Optional[str]) -> set[str]:
    result = set(DEFAULT_IGNORES)
    if raw:
        for token in raw.split(","):
            token = token.strip()
            if token:
                result.add(token)
    return result


def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="document_search.py",
        description="Grep for a term across a workspace — including inside "
                    ".docx, .pptx, .xlsx, and .pdf files. Returns a ranked "
                    "list of files that contain the term.",
    )
    p.add_argument("pattern",
                   help="Regex (or literal string with -F) to search for.")
    p.add_argument("path", nargs="?", default=".",
                   help="Root directory to scan (default: current directory).")
    p.add_argument("--ext", default=None,
                   help="Comma-separated list of extensions to include "
                        "(e.g. '.py,.md,.docx'). Defaults to 60+ built-in types.")
    p.add_argument("-F", "--fixed-string", action="store_true",
                   help="Treat PATTERN as a literal string, not a regex.")
    p.add_argument("-i", "--ignore-case", action="store_true",
                   help="Case-insensitive match.")
    p.add_argument("--show-matches", action="store_true",
                   help="Include matched lines (with line numbers) in the output.")
    p.add_argument("--context", type=int, default=0,
                   help="Lines of context around each match (implies --show-matches).")
    p.add_argument("--content", action="store_true",
                   help="Also include the full extracted file content in JSON output.")
    p.add_argument("--max-bytes", type=int, default=200_000,
                   help="Maximum bytes read per file (default: 200000).")
    p.add_argument("--max-files", type=int, default=0,
                   help="Stop after scanning this many files (0 = no limit).")
    p.add_argument("--ignore", default=None,
                   help="Additional comma-separated directory names to skip.")
    p.add_argument("--include-hidden", action="store_true",
                   help="Include dotfiles and dot-directories.")
    p.add_argument("--format", choices=("text", "json"), default="text",
                   help="Output format (default: text).")
    return p


def main(argv: Optional[list[str]] = None) -> int:
    # Force UTF-8 on stdout/stderr so bilingual install guides render on
    # Windows consoles (default cp1252) and redirected pipes alike.
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")  # type: ignore[attr-defined]
        except (AttributeError, ValueError):
            pass

    args = build_parser().parse_args(argv)

    root = Path(args.path)
    if not root.exists():
        print(f"error: path does not exist: {root}", file=sys.stderr)
        return 2
    if not root.is_dir():
        print(f"error: path is not a directory: {root}", file=sys.stderr)
        return 2

    extensions = parse_extensions(args.ext)
    ignore_dirs = parse_ignores(args.ignore)

    flags = re.IGNORECASE if args.ignore_case else 0
    raw_pattern = re.escape(args.pattern) if args.fixed_string else args.pattern
    try:
        pattern = re.compile(raw_pattern, flags)
    except re.error as e:
        print(f"error: invalid regex: {e}", file=sys.stderr)
        return 2

    show_matches = args.show_matches or args.context > 0

    summary = scan(
        root,
        pattern,
        extensions=extensions,
        ignore_dirs=ignore_dirs,
        max_bytes=args.max_bytes,
        max_files=args.max_files,
        include_hidden=args.include_hidden,
        show_matches=show_matches,
        include_content=args.content,
        context=args.context,
    )

    if args.format == "json":
        print(render_json(summary))
    else:
        print(render_text(summary, show_matches=show_matches))

    # Exit code: 0 if anything matched, 1 if nothing matched (grep-style).
    return 0 if summary.files_matched else 1


if __name__ == "__main__":
    sys.exit(main())
